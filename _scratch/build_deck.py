# -*- coding: utf-8 -*-
"""Build the 'EBR Tooling Overview' deck.
Audience includes an external contractor -> content is method/tool level only:
no client names, product/part numbers, batch-record IDs, order numbers or system IDs."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT = r"c:\Users\carlo\Dev\TechSpecs\EBR Tooling Overview.pptx"

# palette
INK   = RGBColor(0x1D, 0x24, 0x2B)
MUTED = RGBColor(0x5F, 0x6B, 0x776 % 0x100)
MUTED = RGBColor(0x5F, 0x6B, 0x76)
GREEN = RGBColor(0x2F, 0x6D, 0x2F)
LIGHT = RGBColor(0xEE, 0xF4, 0xE2)
RULE  = RGBColor(0xD8, 0xDB, 0xDF)
ACCENT2 = RGBColor(0x2C, 0x7C, 0x9C)
AMBER = RGBColor(0xB2, 0x6A, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def tb(slide, x, y, w, h=0.4, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = align
    return tf

def para(tf, text, size=16, bold=False, color=INK, space_after=8, first=False, level=0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.level = level
    p.space_after = Pt(space_after)
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    r.font.name = "Segoe UI"
    return p

def rule(slide, y=1.28, x=0.6, w=12.13):
    ln = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Emu(9525))
    ln.fill.solid(); ln.fill.fore_color.rgb = RULE
    ln.line.fill.background(); ln.shadow.inherit = False

def header(slide, title, kicker=None):
    if kicker:
        t = tb(slide, 0.6, 0.34, 12.13, 0.3)
        para(t, kicker.upper(), size=11, bold=True, color=GREEN, first=True)
        t2 = tb(slide, 0.6, 0.62, 12.13, 0.6)
        para(t2, title, size=30, bold=True, first=True)
    else:
        t = tb(slide, 0.6, 0.45, 12.13, 0.7)
        para(t, title, size=32, bold=True, first=True)
    rule(slide)

def card(slide, x, y, w, h, title, body_lines, tint=LIGHT, tcolor=GREEN):
    s = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = tint
    s.line.color.rgb = RULE; s.line.width = Pt(0.75); s.shadow.inherit = False
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.22); tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.16); tf.margin_bottom = Inches(0.12)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    para(tf, title, size=15, bold=True, color=tcolor, first=True, space_after=6)
    for ln in body_lines:
        para(tf, ln, size=12, color=INK, space_after=4)
    return s

def arrow(slide, x, y, w=0.42, h=0.3):
    a = slide.shapes.add_shape(33, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = RGBColor(0xB6, 0xBA, 0xC0)
    a.line.fill.background(); a.shadow.inherit = False

def bullets(slide, x, y, w, items, size=16, gap=10):
    tf = tb(slide, x, y, w, 5.0)
    first = True
    for it in items:
        if isinstance(it, tuple):
            txt, lvl = it
        else:
            txt, lvl = it, 0
        para(tf, ("• " if lvl == 0 else "– ") + txt,
             size=size if lvl == 0 else size - 2,
             color=INK if lvl == 0 else MUTED,
             space_after=gap, first=first, level=lvl)
        first = False
    return tf

# ---------------------------------------------------------------- 1 title
s = prs.slides.add_slide(BLANK)
band = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(2.4))
band.fill.solid(); band.fill.fore_color.rgb = LIGHT
band.line.fill.background(); band.shadow.inherit = False
t = tb(s, 0.9, 0.75, 11.5, 1.0)
para(t, "EBR Tooling Overview", size=42, bold=True, first=True)
t = tb(s, 0.9, 1.62, 11.5, 0.5)
para(t, "How we get from a paper batch record to a verified electronic step",
     size=18, color=MUTED, first=True)
t = tb(s, 0.9, 3.0, 11.5, 2.4)
para(t, "Three tools, one pipeline", size=20, bold=True, first=True, space_after=14)
para(t, "1.  Batch Record Scanning  —  read the paper record, design the electronic steps", size=16, space_after=10)
para(t, "2.  Functional / Technical Spec Build  —  turn each step into a build-ready design spec", size=16, space_after=10)
para(t, "3.  EBR Puppeteer  —  drive the live EBR app to explore and verify", size=16, space_after=10)
t = tb(s, 0.9, 6.35, 11.5, 0.4)
para(t, "Tool + method overview — no client, product or system specifics.", size=12, color=MUTED, first=True)

# ---------------------------------------------------------------- 2 pipeline
s = prs.slides.add_slide(BLANK)
header(s, "One pipeline, three tools", "the big picture")
y = 2.0
card(s, 0.6, y, 2.55, 2.3, "Paper batch record",
     ["The client's manufacturing", "directions / MPR.", "", "Prose, tables, sign-offs."],
     tint=RGBColor(0xF4, 0xF6, 0xF8), tcolor=INK)
arrow(s, 3.28, y + 1.0)
card(s, 3.85, y, 2.55, 2.3, "1 · Scanning",
     ["Identify each step,", "reuse vs new,", "draw the mock-ups.", "", "→ mock-ups + EBR sheet"])
arrow(s, 6.53, y + 1.0)
card(s, 7.1, y, 2.55, 2.3, "2 · Spec build",
     ["13-section design spec", "incl. function modules", "+ pseudocode.", "", "→ build-ready .docx"])
arrow(s, 9.78, y + 1.0)
card(s, 10.35, y, 2.4, 2.3, "3 · Puppeteer",
     ["Drive the built step in", "the live EBR app.", "", "→ explored / verified"],
     tint=RGBColor(0xE8, 0xF1, 0xF6), tcolor=ACCENT2)
t = tb(s, 0.6, 4.65, 12.13, 1.8)
para(t, "Why it's one pipeline", size=17, bold=True, first=True, space_after=10)
para(t, "• Each stage produces the input the next stage needs — the mock-up is the spec's picture, the spec is the builder's instruction, Puppeteer proves the result.", size=14, color=MUTED, space_after=8)
para(t, "• Everything is written down in build guides, so the output is consistent regardless of who does the work.", size=14, color=MUTED, space_after=8)
para(t, "• Reuse is the default at every stage — we assemble from a library far more often than we invent.", size=14, color=MUTED)

# ---------------------------------------------------------------- 3 tool 1
s = prs.slides.add_slide(BLANK)
header(s, "Batch Record Scanning", "tool 1 · paper → designed steps")
bullets(s, 0.6, 1.6, 6.0, [
    "Read the batch record end to end and identify what each section actually captures.",
    ("Data vs instruction — does it record values, or just tell the operator to act?", 1),
    ("Recurring patterns — sampling, weighing, timers, sign-offs repeat everywhere.", 1),
    "Decide reuse vs new against the existing step library.",
    ("Most steps already exist — reuse beats rebuilding.", 1),
    ("Only genuinely new behaviour becomes a new step.", 1),
    "Draw a mock-up per step in one of three formats.",
    ("Table · Form · Long-text instruction", 1),
])
card(s, 7.0, 1.6, 5.73, 2.35, "What comes out",
     ["• One mock-up image per step — the visual contract for the build.",
      "• An assembled EBR / PI-sheet: the new electronic record, end to end.",
      "• A crosswalk mapping every step back to its batch-record section."])
card(s, 7.0, 4.2, 5.73, 2.2, "Why it matters",
     ["• The client can compare new vs old side by side, section by section.",
      "• Gaps get caught on paper — before anyone builds anything.",
      "• Reuse decisions are made once, up front."],
     tint=RGBColor(0xF4, 0xF6, 0xF8), tcolor=INK)

# ---------------------------------------------------------------- 4 tool 1 detail
s = prs.slides.add_slide(BLANK)
header(s, "Scanning: the method is written down", "tool 1 · how we keep it consistent")
t = tb(s, 0.6, 1.55, 12.13, 0.5)
para(t, "A build guide encodes the rules so two people scanning the same record land in the same place.", size=15, color=MUTED, first=True)
card(s, 0.6, 2.25, 3.9, 2.0, "Recognition cues",
     ["Read the paper and know what it is:",
      "• a value + a range → a checked entry",
      "• a blank + a signature → a record",
      "• 'per SOP…' with no data → instruction"])
card(s, 4.72, 2.25, 3.9, 2.0, "Hard rules",
     ["Learned the hard way:",
      "• one signature per table line",
      "• times are stamped, never typed",
      "• computed fields are read-only"])
card(s, 8.84, 2.25, 3.89, 2.0, "Reuse library",
     ["The recurring blocks:",
      "• signatures, equipment, materials",
      "• sampling, calculations, timers",
      "• long-text instructions"])
t = tb(s, 0.6, 4.6, 12.13, 2.0)
para(t, "The judgement calls the guide exists to settle", size=17, bold=True, first=True, space_after=10)
para(t, "• Is this one step or three?  Split where the sign-off boundary is — one perform/verify pair = one step.", size=14, color=MUTED, space_after=8)
para(t, "• Does the format follow the label or the data?  The data. A calculation is a form when it's one record, a table when it repeats.", size=14, color=MUTED, space_after=8)
para(t, "• A name match is a candidate, not a fit — we open the existing step and check it actually covers the requirement.", size=14, color=MUTED)

# ---------------------------------------------------------------- 5 tool 2
s = prs.slides.add_slide(BLANK)
header(s, "Functional / Technical Spec Build", "tool 2 · mock-up → build-ready spec")
bullets(s, 0.6, 1.6, 6.2, [
    "Every step gets one design specification, in a fixed 13-section template.",
    ("Purpose, overview, assumptions, validation checks, layout…", 1),
    ("…then the two technical sections that matter most.", 1),
    "Function Module(s) + Pseudocode — the part clients' docs usually lack.",
    ("We design the FM set from the mock-up's controls.", 1),
    ("Then confirm each one against the real system.", 1),
    "The mock-up image is embedded in the spec — developers build to the picture.",
])
card(s, 7.1, 1.6, 5.63, 2.5, "The core discipline",
     ["Reuse an existing function first;",
      "author a new one only when nothing fits.",
      "",
      "And when nothing fits — say so in the spec,",
      "rather than forcing a bad match."],
     tint=RGBColor(0xE8, 0xF1, 0xF6), tcolor=ACCENT2)
card(s, 7.1, 4.35, 5.63, 2.05, "⭐ The lesson that saves rework",
     ["A name match — or even a pattern match —",
      "is a candidate, not a fit.",
      "",
      "Open the real interface and check its",
      "parameters against the step's actual fields."],
     tint=RGBColor(0xFD, 0xF6, 0xE7), tcolor=AMBER)

# ---------------------------------------------------------------- 6 tool 2 detail
s = prs.slides.add_slide(BLANK)
header(s, "Spec build: from controls to functions", "tool 2 · how the technical set is designed")
t = tb(s, 0.6, 1.55, 12.13, 0.4)
para(t, "We read the mock-up's controls and map each to the behaviour it implies:", size=15, color=MUTED, first=True)
rows = [
    ("A ▶ Record button stamping a date/time", "a system-time function — never a typed field"),
    ("A numeric field with a min–max range", "a range validator"),
    ("An entry that fills a read-only field beside it", "an input validation on the entry"),
    ("A repeating table with + Add Row", "a row indexer"),
    ("A per-row signature / a footer sign-off", "signature + audit functions"),
    ("A material line that consumes stock", "a goods-issue process message"),
]
y = 2.15
for i, (a, b) in enumerate(rows):
    tint = RGBColor(0xF7, 0xF8, 0xF9) if i % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
    bg = s.shapes.add_shape(1, Inches(0.6), Inches(y), Inches(12.13), Inches(0.62))
    bg.fill.solid(); bg.fill.fore_color.rgb = tint
    bg.line.color.rgb = RULE; bg.line.width = Pt(0.5); bg.shadow.inherit = False
    tf = tb(s, 0.8, y + 0.13, 5.6)
    para(tf, a, size=13, bold=True, first=True)
    ar = s.shapes.add_shape(33, Inches(6.5), Inches(y + 0.2), Inches(0.35), Inches(0.22))
    ar.fill.solid(); ar.fill.fore_color.rgb = RGBColor(0xB6, 0xBA, 0xC0)
    ar.line.fill.background(); ar.shadow.inherit = False
    tf = tb(s, 7.1, y + 0.13, 5.4)
    para(tf, b, size=13, color=MUTED, first=True)
    y += 0.68
t = tb(s, 0.6, 6.35, 12.13, 0.5)
para(t, "Then every candidate is verified against the live system before it goes in the spec.",
     size=14, bold=True, color=GREEN, first=True)

# ---------------------------------------------------------------- 7 tool 3
s = prs.slides.add_slide(BLANK)
header(s, "EBR Puppeteer", "tool 3 · driving the live EBR app")
bullets(s, 0.6, 1.6, 6.2, [
    "A browser-automation harness that drives the EBR app the way an operator would.",
    ("Launch + log in + navigate to a recipe, and stay open.", 1),
    ("Then attach to that same session to explore or drive it.", 1),
    "Three modes:",
    ("probe / drive — click, type, read, set a cell, sign", 1),
    ("eval — run a query inside the sheet and get structured data back", 1),
    "The workhorse is eval: inspect and measure the real sheet.",
    ("Screenshots are the exception, not the method.", 1),
])
card(s, 7.1, 1.6, 5.63, 2.35, "Why we need it",
     ["• The built step is only 'done' when it behaves correctly in the real app.",
      "• Manual click-through doesn't scale and isn't repeatable.",
      "• It tells us what the app actually does — not what we assumed."],
     tint=RGBColor(0xE8, 0xF1, 0xF6), tcolor=ACCENT2)
card(s, 7.1, 4.2, 5.63, 2.2, "What it gives back",
     ["• Confirms field labels, formats and validation behaviour.",
      "• Surfaces mismatches between the spec/mock-up and reality.",
      "• Turns 'it should work' into 'we drove it'."],
     tint=RGBColor(0xF4, 0xF6, 0xF8), tcolor=INK)

# ---------------------------------------------------------------- 8 tool 3 detail
s = prs.slides.add_slide(BLANK)
header(s, "Puppeteer: why it's harder than it looks", "tool 3 · the realities of the sheet")
t = tb(s, 0.6, 1.55, 12.13, 0.45)
para(t, "The recipe sheet is legacy framework HTML inside an iframe — not a clean modern UI. That shapes everything:",
     size=15, color=MUTED, first=True)
card(s, 0.6, 2.2, 3.9, 2.62, "Fields have no names",
     ["Inputs share one id — no label,",
      "no placeholder, nothing to query.",
      "",
      "So we locate by GEOMETRY: find the",
      "header's position, then the input",
      "beneath it."],
     tint=RGBColor(0xFD, 0xF6, 0xE7), tcolor=AMBER)
card(s, 4.72, 2.2, 3.9, 2.62, "Enter commits",
     ["A value doesn't register until",
      "Enter is pressed — dropdowns too.",
      "",
      "A committed number visibly",
      "reformats; if it didn't reformat,",
      "it didn't commit."],
     tint=RGBColor(0xFD, 0xF6, 0xE7), tcolor=AMBER)
card(s, 8.84, 2.2, 3.89, 2.62, "Validation is visual",
     ["State lives in the field's styling",
      "— an amber border — not in any",
      "message text.",
      "",
      "An invalid value triggers two",
      "separate popups, both of which",
      "must be handled."],
     tint=RGBColor(0xFD, 0xF6, 0xE7), tcolor=AMBER)
t = tb(s, 0.6, 5.05, 12.13, 1.6)
para(t, "Open gaps — worth knowing before you rely on it", size=17, bold=True, first=True, space_after=10)
para(t, "• Setting a cell doesn't yet auto-handle the two validation popups.", size=14, color=MUTED, space_after=7)
para(t, "• Text matching only sees visible text — icon-only buttons need an explicit selector.", size=14, color=MUTED, space_after=7)
para(t, "• Some popups aren't caught by the standard dialog detectors.", size=14, color=MUTED)

# ---------------------------------------------------------------- 9 value
s = prs.slides.add_slide(BLANK)
header(s, "What this adds up to", "why the tooling exists")
card(s, 0.6, 1.7, 3.9, 2.3, "Consistency",
     ["The method is written down, not",
      "carried in someone's head.",
      "",
      "Two people, same record → same",
      "answer."])
card(s, 4.72, 1.7, 3.9, 2.3, "Reuse",
     ["We assemble from a library",
      "instead of inventing.",
      "",
      "Fewer objects to build, test,",
      "and validate."])
card(s, 8.84, 1.7, 3.89, 2.3, "Evidence",
     ["Nothing is 'done' on assumption.",
      "",
      "Functions are checked against the",
      "system; steps are driven in the",
      "real app."],
     tint=RGBColor(0xE8, 0xF1, 0xF6), tcolor=ACCENT2)
t = tb(s, 0.6, 4.4, 12.13, 2.0)
para(t, "The through-line", size=17, bold=True, first=True, space_after=10)
para(t, "Each tool removes a class of error early: scanning catches missed data on paper, the spec catches wrong assumptions before the build, Puppeteer catches behaviour that only shows up in the real app.", size=15, color=MUTED, space_after=9)
para(t, "The expensive mistakes are the ones found late. All three tools exist to find them sooner.", size=15, bold=True, color=INK)

# ---------------------------------------------------------------- 10 close
s = prs.slides.add_slide(BLANK)
band = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(2.0))
band.fill.solid(); band.fill.fore_color.rgb = LIGHT
band.line.fill.background(); band.shadow.inherit = False
t = tb(s, 0.9, 0.62, 11.5, 0.8)
para(t, "Recap", size=36, bold=True, first=True)
t = tb(s, 0.9, 2.5, 11.8, 3.2)
para(t, "1 · Batch Record Scanning — read the paper record, decide reuse vs new, draw the steps.", size=18, bold=True, first=True, space_after=6)
para(t, "     Output: mock-ups + an assembled EBR sheet the client can compare against the old record.", size=14, color=MUTED, space_after=16)
para(t, "2 · Functional / Technical Spec — turn each step into a build-ready design spec.", size=18, bold=True, space_after=6)
para(t, "     Output: a 13-section .docx, including the function modules and pseudocode, built to the mock-up.", size=14, color=MUTED, space_after=16)
para(t, "3 · EBR Puppeteer — drive the built step in the live app.", size=18, bold=True, space_after=6)
para(t, "     Output: evidence it behaves — and early warning where it doesn't.", size=14, color=MUTED)
t = tb(s, 0.9, 6.2, 11.5, 0.5)
para(t, "Questions?", size=22, bold=True, color=GREEN, first=True)

prs.save(OUT)
print("wrote", OUT, "|", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
